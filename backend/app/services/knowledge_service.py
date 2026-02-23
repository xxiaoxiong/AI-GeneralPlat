import os
import re
import math
import json
import sqlite3
import asyncio
from typing import List, Dict, Any, Optional
from pathlib import Path
from collections import Counter

from app.core.config import settings

# SQLite-based hybrid retrieval store
# Combines BM25 + TF-IDF for improved recall

_DB_PATH: Optional[Path] = None


def _get_db_path() -> Path:
    global _DB_PATH
    if _DB_PATH is None:
        persist_dir = Path(settings.CHROMA_PERSIST_DIR)
        persist_dir.mkdir(parents=True, exist_ok=True)
        _DB_PATH = persist_dir / "knowledge.db"
    return _DB_PATH


def _get_conn() -> sqlite3.Connection:
    conn = sqlite3.connect(str(_get_db_path()), check_same_thread=False)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS chunks (
            id TEXT PRIMARY KEY,
            collection TEXT NOT NULL,
            document_id INTEGER NOT NULL,
            filename TEXT NOT NULL,
            chunk_index INTEGER NOT NULL,
            content TEXT NOT NULL
        )
    """)
    conn.execute("CREATE INDEX IF NOT EXISTS idx_collection ON chunks(collection)")
    conn.execute("CREATE INDEX IF NOT EXISTS idx_doc ON chunks(collection, document_id)")
    conn.commit()
    return conn


# BM25 parameters
_BM25_K1 = 1.5
_BM25_B  = 0.75


def _tokenize(text: str) -> List[str]:
    text = text.lower()
    # Chinese character-level + English word-level tokenization
    tokens = re.findall(r'[\u4e00-\u9fff]|[a-z0-9]+', text)
    # Add bigrams for Chinese to improve phrase matching
    chars = re.findall(r'[\u4e00-\u9fff]', text)
    bigrams = [chars[i] + chars[i+1] for i in range(len(chars)-1)]
    return tokens + bigrams


def _bm25_score(
    query_tokens: List[str],
    doc_tokens: List[str],
    doc_freq: Dict[str, int],
    N: int,
    avg_dl: float,
) -> float:
    """BM25 scoring for a single document."""
    if not query_tokens or not doc_tokens:
        return 0.0
    dl = len(doc_tokens)
    tf_map = Counter(doc_tokens)
    score = 0.0
    for qt in query_tokens:
        if qt not in tf_map:
            continue
        tf = tf_map[qt]
        df = doc_freq.get(qt, 0)
        idf = math.log((N - df + 0.5) / (df + 0.5) + 1)
        tf_norm = (tf * (_BM25_K1 + 1)) / (tf + _BM25_K1 * (1 - _BM25_B + _BM25_B * dl / max(avg_dl, 1)))
        score += idf * tf_norm
    return score


def _tfidf_score(query_tokens: List[str], doc_tokens: List[str], doc_freq: Dict[str, int], N: int) -> float:
    """TF-IDF scoring for a single document."""
    if not query_tokens or not doc_tokens:
        return 0.0
    tf_map = Counter(doc_tokens)
    score = 0.0
    for qt in query_tokens:
        tf = tf_map.get(qt, 0) / max(len(doc_tokens), 1)
        df = doc_freq.get(qt, 0)
        idf = math.log((N + 1) / (df + 1)) + 1
        score += tf * idf
    return score


class KnowledgeService:

    @staticmethod
    def _split_text(text: str, chunk_size: int = 500, chunk_overlap: int = 50) -> List[str]:
        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size
            chunks.append(text[start:end])
            start += chunk_size - chunk_overlap
        return [c for c in chunks if c.strip()]

    @staticmethod
    async def extract_text(file_path: str, file_type: str) -> str:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, KnowledgeService._extract_text_sync, file_path, file_type)

    @staticmethod
    def _extract_text_sync(file_path: str, file_type: str) -> str:
        ext = file_type.lower().lstrip(".")
        try:
            if ext == "pdf":
                from pypdf import PdfReader
                reader = PdfReader(file_path)
                return "\n".join(page.extract_text() or "" for page in reader.pages)
            elif ext in ("doc", "docx"):
                from docx import Document
                doc = Document(file_path)
                return "\n".join(p.text for p in doc.paragraphs)
            elif ext in ("xls", "xlsx"):
                import openpyxl
                wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                lines = []
                for sheet in wb.worksheets:
                    lines.append(f"=== Sheet: {sheet.title} ===")
                    for row in sheet.iter_rows(values_only=True):
                        row_vals = [str(c) if c is not None else "" for c in row]
                        if any(v.strip() for v in row_vals):
                            lines.append("\t".join(row_vals))
                return "\n".join(lines)
            elif ext in ("txt", "md", "markdown", "csv"):
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    return f.read()
            elif ext in ("ppt", "pptx"):
                from pptx import Presentation
                prs = Presentation(file_path)
                lines = []
                for i, slide in enumerate(prs.slides, 1):
                    lines.append(f"=== 第{i}页 ===")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            lines.append(shape.text)
                return "\n".join(lines)
            elif ext in ("mp3", "wav", "m4a", "ogg", "flac"):
                return KnowledgeService._transcribe_audio(file_path)
            else:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    return f.read()
        except Exception as e:
            return f"[文档解析失败: {str(e)}]"

    @staticmethod
    def _transcribe_audio(file_path: str) -> str:
        """音频转写：优先使用 faster-whisper，回退到提示信息"""
        try:
            from faster_whisper import WhisperModel
            model = WhisperModel("small", device="cpu", compute_type="int8")
            segments, info = model.transcribe(file_path, beam_size=5)
            text = " ".join(seg.text for seg in segments)
            return f"[音频转写 - 语言: {info.language}]\n{text}"
        except ImportError:
            return "[音频转写需要安装 faster-whisper: pip install faster-whisper]"
        except Exception as e:
            return f"[音频转写失败: {str(e)}]"

    @staticmethod
    async def extract_text_from_url(url: str) -> str:
        """从网页 URL 爬取并提取文本内容"""
        try:
            import httpx
            from html.parser import HTMLParser

            class _TextExtractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self._skip = False
                    self.texts: List[str] = []

                def handle_starttag(self, tag, attrs):
                    if tag in ("script", "style", "nav", "footer", "header"):
                        self._skip = True

                def handle_endtag(self, tag):
                    if tag in ("script", "style", "nav", "footer", "header"):
                        self._skip = False

                def handle_data(self, data):
                    if not self._skip:
                        stripped = data.strip()
                        if stripped:
                            self.texts.append(stripped)

            async with httpx.AsyncClient(timeout=15, verify=False,
                                         headers={"User-Agent": "Mozilla/5.0"}) as client:
                resp = await client.get(url)
                resp.raise_for_status()
                parser = _TextExtractor()
                parser.feed(resp.text)
                text = "\n".join(parser.texts)
                return text[:50000]  # 限制长度
        except Exception as e:
            return f"[网页爬取失败: {str(e)}]"

    @staticmethod
    async def index_document(
        collection_name: str,
        document_id: int,
        filename: str,
        file_path: str,
        file_type: str,
        chunk_size: int = 500,
        chunk_overlap: int = 50,
        embedding_fn=None,
    ) -> int:
        text = await KnowledgeService.extract_text(file_path, file_type)
        chunks = KnowledgeService._split_text(text, chunk_size, chunk_overlap)
        if not chunks:
            return 0

        def _insert():
            conn = _get_conn()
            conn.execute(
                "DELETE FROM chunks WHERE collection=? AND document_id=?",
                (collection_name, document_id)
            )
            rows = [
                (f"doc_{document_id}_chunk_{i}", collection_name, document_id, filename, i, chunk)
                for i, chunk in enumerate(chunks)
            ]
            conn.executemany(
                "INSERT OR REPLACE INTO chunks(id, collection, document_id, filename, chunk_index, content) VALUES(?,?,?,?,?,?)",
                rows
            )
            conn.commit()
            conn.close()

        loop = asyncio.get_event_loop()
        await loop.run_in_executor(None, _insert)
        return len(chunks)

    @staticmethod
    async def search(
        collection_name: str,
        query: str,
        top_k: int = 5,
        score_threshold: float = 0.3,
        embedding_fn=None,
        alpha: float = 0.6,  # BM25 weight; (1-alpha) = TF-IDF weight
    ) -> List[Dict[str, Any]]:
        """Hybrid BM25 + TF-IDF retrieval with score fusion."""
        def _search():
            conn = _get_conn()
            rows = conn.execute(
                "SELECT id, document_id, filename, chunk_index, content FROM chunks WHERE collection=?",
                (collection_name,)
            ).fetchall()
            conn.close()
            return rows

        loop = asyncio.get_event_loop()
        rows = await loop.run_in_executor(None, _search)
        if not rows:
            return []

        query_tokens = _tokenize(query)
        if not query_tokens:
            return []

        all_docs_tokens = [_tokenize(r[4]) for r in rows]
        N = len(all_docs_tokens)

        # Compute doc frequency for both scorers
        doc_freq: Dict[str, int] = {}
        for tokens in all_docs_tokens:
            for t in set(tokens):
                doc_freq[t] = doc_freq.get(t, 0) + 1

        # Average document length for BM25
        avg_dl = sum(len(t) for t in all_docs_tokens) / max(N, 1)

        query_token_set = set(query_tokens)
        scored = []
        for i, row in enumerate(rows):
            doc_tokens = all_docs_tokens[i]
            if not doc_tokens:
                continue

            # Require at least one query token to appear
            hit_tokens = [qt for qt in query_tokens if qt in set(doc_tokens)]
            if not hit_tokens:
                continue

            bm25  = _bm25_score(query_tokens, doc_tokens, doc_freq, N, avg_dl)
            tfidf = _tfidf_score(query_tokens, doc_tokens, doc_freq, N)

            # Hit rate penalty for partial matches
            hit_rate = len(hit_tokens) / len(query_tokens)
            hybrid = (alpha * bm25 + (1 - alpha) * tfidf) * hit_rate

            scored.append((hybrid, bm25, tfidf, row))

        if not scored:
            return []

        scored.sort(key=lambda x: x[0], reverse=True)
        max_score = scored[0][0]

        items = []
        for hybrid, bm25, tfidf, row in scored[:top_k * 2]:  # fetch extra, filter by threshold
            norm_score = round(hybrid / max_score, 4) if max_score > 0 else 0.0
            if norm_score >= score_threshold:
                items.append({
                    "content": row[4],
                    "score": norm_score,
                    "document_id": row[1],
                    "filename": row[2],
                    "chunk_index": row[3],
                    "metadata": {
                        "document_id": row[1],
                        "filename": row[2],
                        "chunk_index": row[3],
                        "bm25_score": round(bm25, 4),
                        "tfidf_score": round(tfidf, 4),
                    },
                })
            if len(items) >= top_k:
                break
        return items

    @staticmethod
    def delete_document_chunks(collection_name: str, document_id: int):
        try:
            conn = _get_conn()
            conn.execute(
                "DELETE FROM chunks WHERE collection=? AND document_id=?",
                (collection_name, document_id)
            )
            conn.commit()
            conn.close()
        except Exception:
            pass

    @staticmethod
    def delete_collection(collection_name: str):
        try:
            conn = _get_conn()
            conn.execute("DELETE FROM chunks WHERE collection=?", (collection_name,))
            conn.commit()
            conn.close()
        except Exception:
            pass
